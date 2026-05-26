#!/usr/bin/env python3
"""
Generate Admin and Learner How-To tutorial PPTX decks
for the AA Plus Policy Training Platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────────────────
BG     = RGBColor(11,  15,  25)   # #0B0F19  dark navy
CARD   = RGBColor(17,  27,  46)   # panel
PURPLE = RGBColor(139, 92,  246)  # #8B5CF6
CYAN   = RGBColor(34,  211, 238)  # #22D3EE
PINK   = RGBColor(236, 72,  153)  # #EC4899
WHITE  = RGBColor(255, 255, 255)
GRAY   = RGBColor(148, 163, 184)  # slate-400
LGRAY  = RGBColor(203, 213, 225)  # slate-300
ORANGE = RGBColor(251, 146, 60)   # orange-400

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(1.5)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width     = line_width
    else:
        s.line.fill.background()
    return s


def add_txt(slide, text, l, t, w, h, size=18, bold=False,
            color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text            = text
    r.font.size       = Pt(size)
    r.font.bold       = bold
    r.font.italic     = italic
    r.font.color.rgb  = color
    r.font.name       = "Calibri"
    return box


def add_multiline(slide, lines, l, t, w, h):
    """
    lines: list of dicts with keys text, size, bold, color, italic, spacing_after
    """
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if ln.get("spacing_after"):
            p.space_after = Pt(ln["spacing_after"])
        if ln.get("spacing_before"):
            p.space_before = Pt(ln["spacing_before"])
        r = p.add_run()
        r.text           = ln.get("text", "")
        r.font.size      = Pt(ln.get("size", 18))
        r.font.bold      = ln.get("bold", False)
        r.font.italic    = ln.get("italic", False)
        r.font.color.rgb = ln.get("color", WHITE)
        r.font.name      = "Calibri"
    return box


# ── Slide templates ───────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle, brand_color=PURPLE):
    """Full-bleed title slide."""
    s = blank(prs)
    set_bg(s, BG)

    # Left accent bar
    add_rect(s, Inches(0), Inches(0), Inches(0.35), H, fill_color=brand_color)

    # Decorative circle (top-right)
    circ = s.shapes.add_shape(9, Inches(10.5), Inches(-1.2), Inches(4.5), Inches(4.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(30, 20, 55)
    circ.line.fill.background()

    # Brand pill
    add_rect(s, Inches(0.7), Inches(1.1), Inches(2.9), Inches(0.42), fill_color=brand_color)
    add_txt(s, "AA PLUS POLICY TRAINING", Inches(0.75), Inches(1.1), Inches(2.8), Inches(0.42),
            size=10, bold=True, color=WHITE)

    # Title
    add_txt(s, title, Inches(0.7), Inches(1.8), Inches(9.5), Inches(1.6),
            size=42, bold=True, color=WHITE)

    # Accent line
    add_rect(s, Inches(0.7), Inches(3.55), Inches(1.6), Inches(0.05), fill_color=brand_color)

    # Subtitle
    add_txt(s, subtitle, Inches(0.7), Inches(3.7), Inches(9), Inches(0.9),
            size=20, color=LGRAY)

    # Bottom bar
    add_rect(s, Inches(0), Inches(6.9), W, Inches(0.6), fill_color=CARD)
    add_txt(s, "training.aaplus.app", Inches(0.7), Inches(6.92), Inches(5), Inches(0.45),
            size=13, color=GRAY)

    return s


def slide_toc(prs, sections, brand_color=PURPLE):
    """Table of contents slide."""
    s = blank(prs)
    set_bg(s, BG)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill_color=brand_color)

    add_txt(s, "What's Inside", Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
            size=32, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(1.3), Inches(1.2), Inches(0.05), fill_color=brand_color)

    col_w   = Inches(5.6)
    x_start = [Inches(0.5), Inches(6.9)]
    y_start = Inches(1.55)
    row_h   = Inches(1.22)

    for i, (num, sec_title, items) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = x_start[col]
        y = y_start + row * row_h

        # Card
        add_rect(s, x, y, col_w, row_h - Inches(0.08), fill_color=CARD)

        # Number badge
        badge = s.shapes.add_shape(9, x + Inches(0.15), y + Inches(0.15),
                                   Inches(0.48), Inches(0.48))
        badge.fill.solid()
        badge.fill.fore_color.rgb = brand_color
        badge.line.fill.background()

        add_txt(s, num, x + Inches(0.15), y + Inches(0.13), Inches(0.48), Inches(0.48),
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_txt(s, sec_title, x + Inches(0.72), y + Inches(0.14), col_w - Inches(0.9),
                Inches(0.42), size=16, bold=True, color=WHITE)

        item_lines = [{"text": f"• {it}", "size": 13, "color": LGRAY, "spacing_after": 1}
                      for it in items]
        add_multiline(s, item_lines, x + Inches(0.72), y + Inches(0.56),
                      col_w - Inches(0.9), Inches(0.65))

    return s


def slide_section(prs, number, title, description="", brand_color=PURPLE):
    """Section divider slide."""
    s = blank(prs)
    set_bg(s, BG)

    # Large background accent shape
    add_rect(s, Inches(6.5), Inches(0), Inches(6.83), H, fill_color=CARD)

    # Vertical accent bar
    add_rect(s, Inches(6.5), Inches(0), Inches(0.22), H, fill_color=brand_color)

    # Big section number
    add_txt(s, number, Inches(0.5), Inches(1.0), Inches(5.5), Inches(2.5),
            size=110, bold=True, color=RGBColor(25, 35, 60), align=PP_ALIGN.LEFT)

    add_txt(s, "SECTION", Inches(0.5), Inches(1.05), Inches(3), Inches(0.5),
            size=13, bold=True, color=brand_color)

    add_txt(s, title, Inches(0.5), Inches(3.0), Inches(5.8), Inches(1.2),
            size=36, bold=True, color=WHITE)

    if description:
        add_rect(s, Inches(0.5), Inches(4.35), Inches(0.9), Inches(0.05), fill_color=brand_color)
        add_txt(s, description, Inches(0.5), Inches(4.5), Inches(5.8), Inches(1.5),
                size=16, color=LGRAY)

    add_txt(s, "AA Plus Policy Training Platform",
            Inches(6.9), Inches(6.4), Inches(5.5), Inches(0.5),
            size=13, color=GRAY)

    return s


def slide_steps(prs, title, steps, note="", brand_color=PURPLE, icon_char="→"):
    """
    Steps slide — title on left panel, numbered steps on right.
    steps: list of (header, detail) or just str
    """
    s = blank(prs)
    set_bg(s, BG)

    # Left panel
    add_rect(s, Inches(0), Inches(0), Inches(4.2), H, fill_color=CARD)
    add_rect(s, Inches(0), Inches(0), Inches(0.22), H, fill_color=brand_color)

    # Icon circle
    circ = s.shapes.add_shape(9, Inches(0.45), Inches(0.45), Inches(0.9), Inches(0.9))
    circ.fill.solid()
    circ.fill.fore_color.rgb = brand_color
    circ.line.fill.background()
    add_txt(s, icon_char, Inches(0.44), Inches(0.45), Inches(0.92), Inches(0.92),
            size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Left panel title
    add_txt(s, title, Inches(0.35), Inches(1.6), Inches(3.65), Inches(3.5),
            size=26, bold=True, color=WHITE)

    add_rect(s, Inches(0.35), Inches(5.3), Inches(1.0), Inches(0.06),
             fill_color=brand_color)

    add_txt(s, "Step-by-step guide", Inches(0.35), Inches(5.5), Inches(3.5), Inches(0.5),
            size=12, color=GRAY, italic=True)

    # Right panel — steps
    right_x = Inches(4.55)
    step_y  = Inches(0.55)
    step_w  = Inches(8.5)

    for i, step in enumerate(steps):
        if isinstance(step, str):
            header, detail = step, ""
        else:
            header, detail = step[0], step[1] if len(step) > 1 else ""

        # Number circle
        cx = right_x
        cy = step_y + i * Inches(1.03)
        num_circ = s.shapes.add_shape(9, cx, cy + Inches(0.08), Inches(0.42), Inches(0.42))
        num_circ.fill.solid()
        num_circ.fill.fore_color.rgb = brand_color
        num_circ.line.fill.background()
        add_txt(s, str(i + 1), cx, cy + Inches(0.06), Inches(0.42), Inches(0.44),
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Connector line between steps
        if i < len(steps) - 1:
            add_rect(s, cx + Inches(0.19), cy + Inches(0.5),
                     Inches(0.04), Inches(0.6),
                     fill_color=RGBColor(50, 60, 90))

        # Step text
        lines = [{"text": header, "size": 16, "bold": True, "color": WHITE}]
        if detail:
            lines.append({"text": detail, "size": 13, "color": LGRAY, "spacing_before": 1})
        add_multiline(s, lines, cx + Inches(0.55), cy + Inches(0.06),
                      step_w - Inches(0.6), Inches(0.95))

    # Note / tip at bottom
    if note:
        add_rect(s, right_x, Inches(6.65), step_w, Inches(0.68), fill_color=CARD)
        add_rect(s, right_x, Inches(6.65), Inches(0.06), Inches(0.68),
                 fill_color=CYAN)
        add_txt(s, f"Tip: {note}", right_x + Inches(0.18), Inches(6.68),
                step_w - Inches(0.25), Inches(0.62), size=13, color=LGRAY, italic=True)

    return s


def slide_summary(prs, title, points, brand_color=PURPLE):
    """Closing / summary slide."""
    s = blank(prs)
    set_bg(s, BG)

    # Top bar
    add_rect(s, Inches(0), Inches(0), W, Inches(0.25), fill_color=brand_color)

    add_txt(s, title, Inches(0.6), Inches(0.5), Inches(12), Inches(1.0),
            size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(5.5), Inches(1.55), Inches(2.33), Inches(0.05),
             fill_color=brand_color)

    col_w = Inches(3.8)
    for i, (icon, heading, body) in enumerate(points):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.28)
        y = Inches(1.85) + row * Inches(2.3)

        add_rect(s, x, y, col_w, Inches(2.1), fill_color=CARD)
        add_rect(s, x, y, col_w, Inches(0.08), fill_color=brand_color)

        add_txt(s, icon, x + Inches(0.2), y + Inches(0.15),
                Inches(0.7), Inches(0.6), size=22, color=brand_color)
        add_txt(s, heading, x + Inches(0.2), y + Inches(0.7),
                col_w - Inches(0.3), Inches(0.5), size=15, bold=True, color=WHITE)
        add_txt(s, body, x + Inches(0.2), y + Inches(1.2),
                col_w - Inches(0.3), Inches(0.8), size=13, color=LGRAY)

    # Footer
    add_rect(s, Inches(0), Inches(6.9), W, Inches(0.6), fill_color=CARD)
    add_txt(s, "Questions? Contact your organisation's admin or visit training.aaplus.app",
            Inches(0), Inches(6.92), W, Inches(0.45),
            size=13, color=GRAY, align=PP_ALIGN.CENTER)

    return s


# ═══════════════════════════════════════════════════════════════════════════════
# ADMIN DECK
# ═══════════════════════════════════════════════════════════════════════════════

def build_admin_deck():
    prs = new_prs()

    # 1 — Title
    slide_title(prs,
                "Admin How-To Guide",
                "A complete reference for managing your organisation's compliance training",
                brand_color=PURPLE)

    # 2 — TOC
    slide_toc(prs, [
        ("01", "Inviting & Managing Users",
         ["Single & bulk user invitations", "Changing roles", "Revoking access"]),
        ("02", "Departments & Modules",
         ["Create departments", "Assign modules", "Set deadlines"]),
        ("03", "Reports & Compliance",
         ["Compliance dashboard", "Audit log export", "Certificate downloads"]),
        ("04", "AI Policy & Org Settings",
         ["AI policy quiz", "Upload logo", "POSH / ICC details"]),
    ], brand_color=PURPLE)

    # ── Section 01 ────────────────────────────────────────────────────────────
    slide_section(prs, "01", "Inviting &\nManaging Users",
                  "Add individuals or bulk-import teams, assign roles, and control access.",
                  brand_color=PURPLE)

    slide_steps(prs, "Invite a Single User",
                [("Go to Admin Dashboard → Users tab",
                  "Click Users in the top navigation of the admin panel."),
                 ("Click the Invite User button",
                  "Found in the top-right corner of the Users tab."),
                 ("Enter the user's name and email address",
                  "Both fields are required to send the invitation."),
                 ("Select a role: Learner or Manager",
                  "Managers can view their team's compliance status."),
                 ("(Optional) Select a department",
                  "Assigns the user to a team and filters their modules."),
                 ("Click Send Invite",
                  "The user receives a setup email with a secure link.")],
                note="Invitations expire after 7 days — resend from the Users tab if needed.",
                brand_color=PURPLE, icon_char="✉")

    slide_steps(prs, "Bulk Import Users\nvia CSV",
                [("Go to Admin Dashboard → Users tab",
                  ""),
                 ("Click Bulk Invite and download the CSV template",
                  "The template has Name, Email, and Department columns."),
                 ("Fill in one user per row and save the file",
                  "Department names must match exactly what's in the platform."),
                 ("Upload the completed CSV file",
                  "Drag and drop or click Browse to select the file."),
                 ("Review the preview list for any errors",
                  "Rows with invalid emails are highlighted in red."),
                 ("Click Send All Invites",
                  "Each user receives an individual invitation email.")],
                note="Use the template — custom column orders will cause import errors.",
                brand_color=PURPLE, icon_char="⬆")

    slide_steps(prs, "Change a User's Role",
                [("Go to Admin Dashboard → Users tab",
                  ""),
                 ("Find the user with the search bar",
                  "Search by name or email address."),
                 ("Click the role badge next to their name",
                  "The badge shows their current role (Learner / Manager)."),
                 ("Select the new role from the dropdown",
                  ""),
                 ("The change takes effect immediately",
                  "No email is sent — the user sees their new access on next login.")],
                note="Admins cannot demote themselves; ask another admin or your SuperAdmin.",
                brand_color=PURPLE, icon_char="↕")

    slide_steps(prs, "Revoke User Access",
                [("Go to Admin Dashboard → Users tab",
                  ""),
                 ("Find the user and click the ··· menu",
                  "The three-dot menu appears at the right of each user row."),
                 ("Select Revoke Access",
                  ""),
                 ("Confirm the action in the dialog",
                  "Type CONFIRM if prompted to prevent accidental revocation."),
                 ("The user is immediately signed out",
                  "They can no longer log in or access any training content.")],
                note="Revoked users can be re-invited at any time with the same email.",
                brand_color=PURPLE, icon_char="✕")

    # ── Section 02 ────────────────────────────────────────────────────────────
    slide_section(prs, "02", "Departments &\nModule Assignment",
                  "Structure your organisation into teams and control which training is required for each.",
                  brand_color=CYAN)

    slide_steps(prs, "Create a Department",
                [("Go to Admin Dashboard → Departments tab",
                  ""),
                 ("Click Add Department",
                  ""),
                 ("Enter a department name (e.g. Finance, Legal, HR)",
                  "Names must be unique within your organisation."),
                 ("Click Save",
                  "The department is immediately available for user assignment."),
                 ("(Optional) Enable RPT Simulator for the department",
                  "Toggle the simulator switch on the department card.")],
                note="You can delete departments only if they have no assigned users.",
                brand_color=CYAN, icon_char="＋")

    slide_steps(prs, "Assign Users to\na Department",
                [("Go to Admin Dashboard → Users tab",
                  ""),
                 ("Select one or more users using the checkboxes",
                  "Tick the header checkbox to select the entire visible list."),
                 ("Click Assign Department in the bulk-action toolbar",
                  "The toolbar appears at the top when users are selected."),
                 ("Choose the target department from the dropdown",
                  ""),
                 ("Click Confirm",
                  "Users are reassigned instantly; their module list updates automatically.")],
                note="A user can only belong to one department at a time.",
                brand_color=CYAN, icon_char="⊕")

    slide_steps(prs, "Enable a Module\nfor Your Organisation",
                [("Go to Admin Dashboard → Modules & Compliance tab",
                  ""),
                 ("Locate the policy module you want to activate",
                  "Modules include Anti-Corruption, POSH, Data Privacy, Cybersecurity, and more."),
                 ("Toggle the Enable switch to ON",
                  ""),
                 ("Choose scope: All Employees or specific departments",
                  ""),
                 ("Click Save",
                  "The module appears on the dashboards of assigned users immediately.")],
                note="Disabling a module hides it but preserves all existing completion records.",
                brand_color=CYAN, icon_char="◎")

    slide_steps(prs, "Set a Module Deadline",
                [("Go to Admin Dashboard → Modules & Compliance tab",
                  ""),
                 ("Click the calendar icon next to the active module",
                  "Only enabled modules can have deadlines."),
                 ("Pick a due date from the date picker",
                  ""),
                 ("Click Save",
                  "The deadline is displayed on every assigned user's dashboard."),
                 ("Overdue warnings appear automatically",
                  "A red banner is shown to any user who misses the deadline.")],
                note="Set deadlines at least 2 weeks ahead to give users enough time.",
                brand_color=CYAN, icon_char="📅")

    # ── Section 03 ────────────────────────────────────────────────────────────
    slide_section(prs, "03", "Reports &\nCompliance Tracking",
                  "Monitor training progress, export audit records, and download certificates.",
                  brand_color=PINK)

    slide_steps(prs, "View the Compliance\nDashboard",
                [("Go to Admin Dashboard → Reports tab",
                  ""),
                 ("View the compliance overview chart",
                  "Shows % of employees who have completed all assigned modules."),
                 ("Use the date range filter to refine the view",
                  "Compare progress across different time periods."),
                 ("Click a department bar to drill down",
                  "See completion status for every individual in that team."),
                 ("Switch to the Leaderboard tab",
                  "View top-performing employees ranked by completion score.")],
                note="The compliance % only counts modules that are marked as required.",
                brand_color=PINK, icon_char="📊")

    slide_steps(prs, "Export the Audit Log",
                [("Go to Admin Dashboard → Reports tab → Audit Log",
                  ""),
                 ("Set the date range filter (From / To)",
                  "Leave blank to export the full history."),
                 ("Click Export CSV",
                  "The file downloads to your computer immediately."),
                 ("Open in Excel or Google Sheets",
                  "The log includes logins, quiz attempts, completions, and admin actions.")],
                note="Audit logs are available for the past 12 months.",
                brand_color=PINK, icon_char="📋")

    slide_steps(prs, "Download Certificates",
                [("Go to Admin Dashboard → Reports tab",
                  ""),
                 ("Click Bulk Download Certificates",
                  ""),
                 ("Select the module and date range",
                  "Filter by completion date to get a specific cohort's certificates."),
                 ("Click Export",
                  "A ZIP file of individual PDF certificates is downloaded."),
                 ("To view a single certificate:",
                  "Find the user in the Users tab → click View Certificate.")],
                note="Certificates include your organisation logo — upload it in Settings first.",
                brand_color=PINK, icon_char="🏆")

    # ── Section 04 ────────────────────────────────────────────────────────────
    slide_section(prs, "04", "AI Policy &\nOrg Settings",
                  "Configure your AI use policy tier, upload branding, and set up POSH committee details.",
                  brand_color=ORANGE)

    slide_steps(prs, "Run the AI Policy Quiz",
                [("Go to Admin Dashboard → AI Policy tab",
                  ""),
                 ("Read the four policy assessment questions carefully",
                  "Questions cover how your organisation uses AI in workflows."),
                 ("Answer each question based on your organisation's AI usage",
                  "Examples: invoice processing, legal drafting, anonymisation."),
                 ("The platform assigns an AI tier automatically:",
                  "Tier 1: Restricted Use  |  Tier 2: Standard  |  Tier 3: Responsible Use"),
                 ("The matching AI training module is assigned to all users",
                  "Users see the new module on their dashboard immediately.")],
                note="You can re-run the quiz if your organisation's AI practices change.",
                brand_color=ORANGE, icon_char="🤖")

    slide_steps(prs, "Upload Your\nOrganisation Logo",
                [("Go to Admin → Settings (gear icon in the sidebar)",
                  ""),
                 ("Scroll to the Organisation Logo section",
                  ""),
                 ("Click Upload Logo",
                  "Supported formats: PNG, JPG — maximum file size 2 MB."),
                 ("Preview the logo in the thumbnail",
                  "Ensure it looks clear at small sizes."),
                 ("Click Save",
                  "The logo appears on all new and existing user certificates.")],
                note="A square or landscape logo with a transparent background looks best.",
                brand_color=ORANGE, icon_char="🖼")

    slide_steps(prs, "Configure POSH /\nICC Details",
                [("Go to Admin → POSH Policy (via the sidebar)",
                  ""),
                 ("Enter your company's registered address and helpline number",
                  ""),
                 ("Enter the Presiding Officer's name and contact details",
                  "The Presiding Officer chairs the Internal Complaints Committee."),
                 ("Add all IC (Internal Committee) members",
                  "Include names and designations for each member."),
                 ("Click Save",
                  "Details are embedded in the POSH policy module content.")],
                note="Update these details immediately if your ICC membership changes.",
                brand_color=ORANGE, icon_char="⚖")

    # ── Summary ────────────────────────────────────────────────────────────────
    slide_summary(prs, "You're All Set!",
                  [("✉", "Invite & Manage",
                    "Single invites, bulk CSV imports, role changes, and access revocation."),
                   ("🏢", "Departments & Modules",
                    "Create teams, enable required training, and set completion deadlines."),
                   ("📊", "Reports",
                    "Monitor compliance, export audit logs, and download certificates."),
                   ("⚙", "Settings",
                    "AI policy tier, organisation logo, and POSH/ICC configuration."),
                   ("❓", "Need help?",
                    "Contact support via training.aaplus.app or your SuperAdmin."),
                   ("🔄", "Keep it updated",
                    "Re-run the AI policy quiz and update ICC details whenever they change.")],
                  brand_color=PURPLE)

    prs.save(r"c:\Users\pc\Desktop\PolicyTraining\AA_Plus_Admin_How-To_Guide.pptx")
    print("Saved: AA_Plus_Admin_How-To_Guide.pptx")


# ═══════════════════════════════════════════════════════════════════════════════
# LEARNER DECK
# ═══════════════════════════════════════════════════════════════════════════════

def build_learner_deck():
    prs = new_prs()

    # 1 — Title
    slide_title(prs,
                "Learner How-To Guide",
                "Everything you need to complete your compliance training with confidence",
                brand_color=CYAN)

    # 2 — TOC
    slide_toc(prs, [
        ("01", "Getting Started",
         ["First login & password setup", "Understanding your dashboard"]),
        ("02", "Completing Training",
         ["Starting a module", "Navigating slides", "Taking & passing the quiz"]),
        ("03", "Your Certificates",
         ["View & print certificates", "Share on LinkedIn"]),
        ("04", "Deadlines & Recertification",
         ["Understanding deadlines", "Overdue warnings", "Annual renewal"]),
    ], brand_color=CYAN)

    # ── Section 01 ────────────────────────────────────────────────────────────
    slide_section(prs, "01", "Getting Started",
                  "Log in for the first time and find your way around the dashboard.",
                  brand_color=CYAN)

    slide_steps(prs, "Log In for the\nFirst Time",
                [("Check your email for an invitation from AA Plus",
                  "Subject: 'You've been invited to AA Plus Policy Training'"),
                 ("Click the Set Up Your Account link in the email",
                  "The link is valid for 7 days — contact your admin if it has expired."),
                 ("You are taken to the password creation page",
                  ""),
                 ("Enter a strong password and confirm it",
                  "Minimum 8 characters — mix letters, numbers, and symbols."),
                 ("Click Save Password",
                  "You are logged in automatically and taken to your training dashboard.")],
                note="Bookmark training.aaplus.app for quick access in future.",
                brand_color=CYAN, icon_char="🔑")

    slide_steps(prs, "Navigate Your\nDashboard",
                [("The large ring shows your overall completion %",
                  "It fills as you complete each assigned module."),
                 ("Module cards list all your assigned training topics",
                  "Each card shows the policy name, status, and deadline."),
                 ("Status labels: Not Started / In Progress / Completed ✓",
                  ""),
                 ("A red banner appears if any module is overdue",
                  "Complete overdue modules as soon as possible."),
                 ("Click any module card to start or resume training",
                  "Your progress is saved — you can always pick up where you left off.")],
                note="Your dashboard updates in real time as you complete each module.",
                brand_color=CYAN, icon_char="🏠")

    # ── Section 02 ────────────────────────────────────────────────────────────
    slide_section(prs, "02", "Completing\nTraining",
                  "Work through training slides and pass the assessment quiz to earn your certificate.",
                  brand_color=PURPLE)

    slide_steps(prs, "Start a Module",
                [("From your dashboard, click on a module card",
                  ""),
                 ("Read the module overview and learning objectives",
                  "This gives you a preview of what the training covers."),
                 ("Click Start Training to begin",
                  ""),
                 ("Your progress is saved automatically",
                  "You can close the browser at any time and resume later."),
                 ("The module reopens exactly where you left off",
                  "Look for the Continue button on the dashboard card.")],
                note="Read every slide carefully — quiz questions are drawn from the content.",
                brand_color=PURPLE, icon_char="▶")

    slide_steps(prs, "Navigate Training\nSlides",
                [("Click Next → to advance through slides",
                  ""),
                 ("Click ← Back to review a previous slide",
                  "You can revisit slides as many times as needed."),
                 ("The progress bar at the top shows how far you've come",
                  "E.g. Slide 6 of 14"),
                 ("Some slides have interactive elements",
                  "Read them carefully before proceeding — they reinforce key points."),
                 ("After the last slide, the quiz becomes available",
                  "Click Take Quiz to start the assessment.")],
                note="There is no time limit on slides — take your time to read thoroughly.",
                brand_color=PURPLE, icon_char="📖")

    slide_steps(prs, "Take the Quiz",
                [("After the last slide, click Take Quiz",
                  ""),
                 ("Answer all 10 multiple-choice questions",
                  "Only one answer is correct per question — choose carefully."),
                 ("You need 70% (7 out of 10) to pass",
                  ""),
                 ("Questions are randomly selected from a larger question bank",
                  "Each attempt presents a different set of questions."),
                 ("Your score is shown immediately after submission",
                  "Pass: your certificate is generated automatically.")],
                note="Re-read the slides if you're unsure — there is no penalty for reviewing.",
                brand_color=PURPLE, icon_char="✏")

    slide_steps(prs, "Retake the Quiz\n(If You Don't Pass)",
                [("If your score is below 70%, click Retake Quiz",
                  ""),
                 ("A new randomised set of 10 questions is generated",
                  "The question bank has more than 10 questions, so expect variety."),
                 ("Review the module slides before retrying",
                  "Focus on topics where you felt less confident."),
                 ("There is no limit on retakes",
                  "Keep going — you'll get there!"),
                 ("Once you pass, your certificate is issued immediately",
                  "It appears on your dashboard and in the certificate view.")],
                note="Most learners pass within 1–2 attempts after reviewing the slides.",
                brand_color=PURPLE, icon_char="🔄")

    # ── Section 03 ────────────────────────────────────────────────────────────
    slide_section(prs, "03", "Your Certificates",
                  "View, print, and share your proof of completion.",
                  brand_color=PINK)

    slide_steps(prs, "View Your Certificate",
                [("After passing the quiz, your certificate is generated automatically",
                  ""),
                 ("Click View Certificate on the result screen",
                  "Or find it via your dashboard — click the module card then View Certificate."),
                 ("The certificate shows your name, module, score, and date",
                  "Your organisation's logo appears in the top corner."),
                 ("Certificates are valid for 12 months from the date of issue",
                  "The expiry date is shown at the bottom of the certificate."),
                 ("All your certificates are stored in your dashboard",
                  "Access them any time — they never disappear from your account.")],
                note="If your name appears incorrectly, contact your admin to update your profile.",
                brand_color=PINK, icon_char="🏆")

    slide_steps(prs, "Print Your Certificate",
                [("Open your certificate (from dashboard or result screen)",
                  ""),
                 ("Click the Print button in the top-right corner",
                  ""),
                 ("The browser print dialog opens",
                  ""),
                 ("Select your printer, or choose Save as PDF",
                  "Saving as PDF lets you email or archive it."),
                 ("The certificate is formatted for A4 / Letter paper",
                  "Choose portrait orientation for the best result.")],
                note="Use 'Save as PDF' to keep a digital copy without printing.",
                brand_color=PINK, icon_char="🖨")

    slide_steps(prs, "Share on LinkedIn",
                [("Open your certificate page",
                  ""),
                 ("Click the Share on LinkedIn button",
                  ""),
                 ("A LinkedIn post is pre-filled with your certification details",
                  "Module name, score, and organisation are included automatically."),
                 ("Customise the post caption if you wish",
                  ""),
                 ("Click Post",
                  "Your network can see your compliance achievement!")],
                note="You may need to log in to LinkedIn if you are not already signed in.",
                brand_color=PINK, icon_char="🔗")

    # ── Section 04 ────────────────────────────────────────────────────────────
    slide_section(prs, "04", "Deadlines &\nRecertification",
                  "Stay on top of due dates and renew your certifications every year.",
                  brand_color=ORANGE)

    slide_steps(prs, "Understanding\nModule Deadlines",
                [("Deadlines are set by your organisation's admin",
                  ""),
                 ("Your dashboard shows the due date on each module card",
                  "E.g. 'Due: 30 June 2026'"),
                 ("A countdown appears as the deadline approaches",
                  "You'll see warnings like '5 days remaining'."),
                 ("Complete the module before the deadline",
                  "No further action is needed once you've passed the quiz."),
                 ("If you need more time, contact your manager or admin",
                  "Only admins can extend deadlines.")],
                note="Start modules early — quizzes may require a couple of attempts.",
                brand_color=ORANGE, icon_char="📅")

    slide_steps(prs, "Overdue Warnings",
                [("If a module is past its deadline, a red Overdue banner appears",
                  "The banner is shown at the top of your dashboard."),
                 ("The overdue module card is highlighted with a warning border",
                  ""),
                 ("Complete the overdue module as soon as possible",
                  "Your compliance status remains incomplete until you pass."),
                 ("Contact your manager or admin if you need an extension",
                  "Admins can adjust the deadline from their dashboard."),
                 ("Your compliance status updates instantly once you pass",
                  "The overdue warning disappears and your completion % increases.")],
                note="Overdue status is visible to your manager — act quickly.",
                brand_color=ORANGE, icon_char="⚠")

    slide_steps(prs, "Recertification —\nAnnual Renewal",
                [("Certificates are valid for 12 months from completion",
                  ""),
                 ("30 days before expiry, a warning appears on your dashboard",
                  "The module card shows 'Expiring Soon'."),
                 ("Click Recertify on the module card",
                  ""),
                 ("Complete the training slides and quiz again",
                  "The content may be updated — read it carefully even if familiar."),
                 ("Pass the quiz to receive a new certificate",
                  "The new certificate resets your 12-month validity period.")],
                note="Recertification keeps your compliance record continuous — don't let it lapse.",
                brand_color=ORANGE, icon_char="♻")

    # ── Summary ────────────────────────────────────────────────────────────────
    slide_summary(prs, "You're Ready to Learn!",
                  [("🔑", "First Login",
                    "Check your invitation email, set your password, and log in."),
                   ("📖", "Complete Modules",
                    "Read every slide, then take the quiz. 70% to pass. Unlimited retakes."),
                   ("🏆", "Certificates",
                    "Issued automatically on passing. Valid for 12 months."),
                   ("🔗", "Share on LinkedIn",
                    "Celebrate your achievement with your professional network."),
                   ("📅", "Watch Deadlines",
                    "Complete modules before the due date shown on your dashboard."),
                   ("♻", "Recertify Annually",
                    "Renew each certificate within 30 days of the expiry warning.")],
                  brand_color=CYAN)

    prs.save(r"c:\Users\pc\Desktop\PolicyTraining\AA_Plus_Learner_How-To_Guide.pptx")
    print("Saved: AA_Plus_Learner_How-To_Guide.pptx")


# ── Entry point ───────────────────────────────────────────────────────────────
if __name__ == "__main__":
    build_admin_deck()
    build_learner_deck()
    print("Done. Both decks are ready.")
