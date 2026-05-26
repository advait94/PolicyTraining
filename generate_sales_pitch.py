#!/usr/bin/env python3
"""
AA Plus Policy Training Platform — Sales Pitch Deck v2
Carefully bounded layouts: every element verified within 13.333 x 7.5 inches.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Dimensions ────────────────────────────────────────────────────────────────
W   = 13.333
H   = 7.5

# ── Margins & layout grid ────────────────────────────────────────────────────
MX          = 0.45          # left/right margin
MY          = 0.35          # top margin for chip
CW          = W - 2 * MX   # 12.433"  usable width
FOOTER_Y    = 6.95          # footer top  (0.55" footer to bottom)
HEADER_END  = 2.15          # content starts after chip+title+rule+subtitle

# ── Palette ───────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0B, 0x0F, 0x19)
SURFACE  = RGBColor(0x15, 0x1A, 0x29)
SURF2    = RGBColor(0x1E, 0x25, 0x38)
PURPLE   = RGBColor(0xA8, 0x55, 0xF7)
PINK     = RGBColor(0xEC, 0x48, 0x99)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)
VIOLET   = RGBColor(0xC0, 0x84, 0xFC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xCB, 0xD5, 0xE1)
GRAY     = RGBColor(0x94, 0xA3, 0xB8)
DGRAY    = RGBColor(0x47, 0x55, 0x69)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)

# ── Primitives ────────────────────────────────────────────────────────────────

def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def box(slide, x, y, w, h, fill=SURFACE, border=None, bw=1.5):
    """Filled rectangle, optionally bordered. All coords in inches."""
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(bw)
    else:
        sh.line.fill.background()
    return sh


def rule(slide, x, y, w, color=PURPLE, h=0.055):
    box(slide, x, y, w, h, fill=color)


def txt(slide, text, x, y, w, h,
        size=11, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color


def chip(slide, text, x, y, w, h=0.30, fill=PURPLE, color=WHITE, size=9):
    box(slide, x, y, w, h, fill=fill)
    txt(slide, text, x + 0.10, y + 0.02, w - 0.20, h - 0.04,
        size=size, bold=True, color=color)


def footer(slide, note="AA Plus Policy Training Platform  |  training.aaplus.app"):
    box(slide, 0, FOOTER_Y, W, H - FOOTER_Y, fill=SURFACE)
    txt(slide, note, MX, FOOTER_Y + 0.10, CW, 0.35,
        size=9, color=DGRAY, align=PP_ALIGN.CENTER)


def header(slide, badge, title, badge_color=PURPLE, subtitle=None, tsize=30):
    """Chip at 0.35, title at 0.80, rule at 1.57, optional subtitle at 1.66."""
    cw = min(len(badge) * 0.115 + 0.55, CW)
    chip(slide, badge, MX, MY, cw, 0.30, fill=badge_color)
    txt(slide, title, MX, 0.80, CW, 0.72,
        size=tsize, bold=True, color=WHITE)
    rule(slide, MX, 1.57, 3.80, color=badge_color)
    if subtitle:
        txt(slide, subtitle, MX, 1.66, CW, 0.38, size=12, color=LGRAY)


# ── Compound components ───────────────────────────────────────────────────────

def info_card(slide, title, bullets, x, y, w, h, accent=PURPLE, bg_color=SURFACE):
    """
    Card layout (all relative to top-left corner of card):
      0.00  rule (accent, 0.055h)
      0.10  title (bold, 13pt, 0.38h)
      0.52  divider (thin, DGRAY)
      0.60  bullets start, 0.38 pitch
    """
    box(slide, x, y, w, h, fill=bg_color)
    rule(slide, x, y, w, color=accent)
    txt(slide, title, x + 0.14, y + 0.10, w - 0.28, 0.38,
        size=13, bold=True, color=WHITE)
    rule(slide, x + 0.14, y + 0.52, w - 0.28, color=DGRAY, h=0.03)
    for i, b in enumerate(bullets):
        txt(slide, f"▸  {b}",
            x + 0.14, y + 0.60 + i * 0.38, w - 0.28, 0.35,
            size=10, color=LGRAY)


def metric(slide, value, label, x, y, w, h, vc=PURPLE):
    box(slide, x, y, w, h, fill=SURF2)
    rule(slide, x, y, w, color=vc)
    txt(slide, value, x, y + 0.10, w, 0.65,
        size=34, bold=True, color=vc, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + 0.78, w, 0.38,
        size=10, color=LGRAY, align=PP_ALIGN.CENTER)


def step(slide, num, heading, body, x, y, w, h=0.88, accent=PURPLE):
    box(slide, x, y, w, h, fill=SURFACE)
    box(slide, x, y, 0.52, h, fill=accent)
    txt(slide, str(num), x, y + (h - 0.38) / 2, 0.52, 0.38,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, heading, x + 0.62, y + 0.06, w - 0.72, 0.32,
        size=11, bold=True, color=WHITE)
    txt(slide, body, x + 0.62, y + 0.42, w - 0.72, 0.38,
        size=9, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# 01  COVER
# ══════════════════════════════════════════════════════════════════════════════
def s_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)

    # Left accent bar  0..0.08 x 0..7.5
    box(slide, 0, 0, 0.08, H, fill=PURPLE)

    # Decorative top-right block  9.5..13.333 x 0..3.4
    box(slide, 9.50, 0, W - 9.50, 3.40, fill=SURF2)

    # Brand chip
    chip(slide, "AA PLUS POLICY TRAINING PLATFORM", 0.28, 0.35, 5.50, 0.30)

    # Headline
    txt(slide, "India's First\nCompliance Training\nPlatform",
        0.28, 0.82, 8.60, 3.00, size=46, bold=True, color=WHITE)

    # Pink rule
    rule(slide, 0.28, 3.88, 3.00, color=PINK)

    # Sub-headline
    txt(slide, "11+ regulatory modules  •  Verifiable certificates\n"
               "Audit-ready reporting  •  Built for India.",
        0.28, 3.98, 8.60, 0.78, size=15, color=LGRAY)

    # URL chip
    chip(slide, "training.aaplus.app", 0.28, 5.00, 2.90, 0.38,
         fill=SURF2, color=CYAN, size=12)

    # Stat cards  x=9.70  w=3.20  x+w=12.90 ✓
    metric(slide, "11+",  "Compliance Modules",   9.70, 1.00, 3.20, 1.25, vc=PURPLE)
    metric(slide, "QR",   "Verified Certificates", 9.70, 2.42, 3.20, 1.25, vc=CYAN)
    metric(slide, "100%", "India-Native Content",  9.70, 3.84, 3.20, 1.25, vc=PINK)

    # Footer
    box(slide, 0, 7.00, W, 0.50, fill=SURFACE)
    txt(slide, "Confidential — For Discussion Purposes Only",
        MX, 7.10, CW / 2, 0.30, size=9, color=DGRAY)
    txt(slide, "training.aaplus.app",
        W / 2, 7.10, CW / 2, 0.30, size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# 02  PROBLEM  (2 x 2 grid)
# CW2=6.08  CH2=2.30  gaps auto
# ══════════════════════════════════════════════════════════════════════════════
def s_problem(prs):
    CW2, CH2, GAP_Y = 6.08, 2.30, 0.14
    GAP_X = CW - 2 * CW2          # 0.273"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "THE PROBLEM",
           "Compliance Training in India is Broken",
           badge_color=PINK,
           subtitle="HR teams pay too much, cover too few regulations, and have nothing to show auditors.")

    problems = [
        (PINK,   "Too Expensive",
         ["Global platforms: SAP Litmos, KnowBe4",
          "Rs 10 lakh+ per year for 500 users",
          "Indian SMBs simply can't afford them"]),
        (ORANGE, "Too Narrow",
         ["POSH Check covers only 1 regulation",
          "No DPDP 2023, Cybersecurity, or AI Policy",
          "Multiple vendors = multiple invoices"]),
        (CYAN,   "No Content Included",
         ["TalentLMS / Disprz: bring your own content",
          "Admin must build all 11 regulation courses",
          "Months of work before first learner logs in"]),
        (PURPLE, "No Audit Trail",
         ["HR tracks completions in spreadsheets",
          "No QR-verified digital certificate",
          "Fails SEBI / DPDP regulatory inspection"]),
    ]
    for i, (ac, title, buls) in enumerate(problems):
        col, row = i % 2, i // 2
        x = MX + col * (CW2 + GAP_X)
        y = HEADER_END + row * (CH2 + GAP_Y)
        info_card(slide, title, buls, x, y, CW2, CH2, accent=ac)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 03  SOLUTION  (4 equal columns)
# CW4 ≈ 3.00  CH4 = FOOTER_Y - HEADER_END - 0.05 = 4.75
# ══════════════════════════════════════════════════════════════════════════════
def s_solution(prs):
    GAP  = 0.14
    CW4  = (CW - 3 * GAP) / 4     # 2.998"
    CH4  = FOOTER_Y - HEADER_END - 0.05   # 4.75"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "THE SOLUTION",
           "One Platform. All Regulations. Audit-Ready.",
           subtitle="India-specific compliance content + automated certificates + full audit trails.")

    pillars = [
        (PURPLE, "Content Included",
         ["11+ pre-built regulatory modules",
          "POSH, DPDP 2023, CERT-In & more",
          "Updated for Indian law",
          "No content add-ons to buy",
          "Slides + quizzes ready Day 1"]),
        (CYAN,   "Proof of Compliance",
         ["Auto-generated PDF certificates",
          "QR code verification on every cert",
          "Bulk certificate export (ZIP)",
          "12-month validity + auto-reminders",
          "LinkedIn sharing built in"]),
        (PINK,   "Admin Control",
         ["Department-level module assignment",
          "Automated email nudges",
          "Excel export + full audit log",
          "Module deadlines + overdue alerts",
          "Slack / MS Teams webhooks"]),
        (ORANGE, "India-First",
         ["INR pricing + GST invoices",
          "Razorpay / UPI AutoPay billing",
          "POSH Wizard with ICC details",
          "DPDP Act 2023 module",
          "AI Policy 3-tier framework"]),
    ]
    for i, (ac, title, buls) in enumerate(pillars):
        x = MX + i * (CW4 + GAP)
        info_card(slide, title, buls, x, HEADER_END, CW4, CH4, accent=ac)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 04  MODULES  (3 columns × 4 rows)
# COL_W=4.03  ROW_H=0.78
# ══════════════════════════════════════════════════════════════════════════════
def s_modules(prs):
    GAP_X = 0.18
    COL_W = (CW - 2 * GAP_X) / 3   # 4.031"
    ROW_H = 0.78
    GAP_Y = 0.08

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "COMPLIANCE MODULES",
           "11+ India-Specific Regulatory Modules",
           subtitle="Every module included in all paid plans — no add-on content fees, no surprises.")

    modules = [
        (PURPLE, "POSH Act 2013",
         "ICC setup, complaint process, POSH Wizard"),
        (CYAN,   "DPDP Act 2023",
         "Digital Personal Data Protection — consent, notices, breach"),
        (PINK,   "CERT-In Cybersecurity",
         "MeitY / CERT-In incident reporting & hygiene"),
        (ORANGE, "Anti-Corruption",
         "FCPA + Prevention of Corruption Act India"),
        (GREEN,  "Code of Conduct",
         "Ethics, conflict of interest, gifts & hospitality"),
        (VIOLET, "AI Policy Tiers",
         "3-tier AI governance: Restricted / Standard / Responsible"),
        (PURPLE, "Data Privacy (General)",
         "GDPR basics, data handling, employee rights"),
        (CYAN,   "Information Security",
         "Passwords, device security, phishing awareness"),
        (PINK,   "Whistleblower Policy",
         "Vigil mechanism, protected disclosure, anti-retaliation"),
        (ORANGE, "Health, Safety & Environment",
         "Factory Act, fire safety, ergonomics basics"),
        (GREEN,  "Company Law + RPT",
         "RPT Simulator — Related Party Transaction training"),
    ]
    for i, (ac, title, desc) in enumerate(modules):
        col, row = i % 3, i // 3
        x = MX + col * (COL_W + GAP_X)
        y = HEADER_END + row * (ROW_H + GAP_Y)
        box(slide, x, y, COL_W, ROW_H, fill=SURFACE)
        rule(slide, x, y, COL_W, color=ac)
        txt(slide, title, x + 0.14, y + 0.10, COL_W - 0.28, 0.32,
            size=12, bold=True)
        txt(slide, desc,  x + 0.14, y + 0.44, COL_W - 0.28, 0.28,
            size=9, color=GRAY)

    # Bottom note bar  y = HEADER_END + 4*(ROW_H+GAP_Y) + 0.08 = 5.71
    ny = HEADER_END + 4 * (ROW_H + GAP_Y) + 0.08
    box(slide, MX, ny, CW, 0.40, fill=SURF2)
    rule(slide, MX, ny, CW, color=PURPLE)
    txt(slide, "  POSH-Only Package available at Rs 25,000/year — undercuts POSH Check (Rs 55,000) while showcasing the full platform.",
        MX + 0.14, ny + 0.06, CW - 0.28, 0.28, size=10, color=LGRAY)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 05  HOW IT WORKS  (2 columns of 3 steps each)
# col_w = (CW - 0.28) / 2 = 6.077"
# Steps: y = HEADER_END + 0.44 + i*1.02, h=0.88
# ══════════════════════════════════════════════════════════════════════════════
def s_how(prs):
    COL_W = (CW - 0.28) / 2   # 6.077"
    RX    = MX + COL_W + 0.28  # right col x = 6.807"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "HOW IT WORKS",
           "Up and Running in Minutes",
           badge_color=CYAN,
           subtitle="Three steps for admins. Three steps for learners. Average first certificate: under 45 minutes.")

    # Column labels
    txt(slide, "For Admins",   MX,  HEADER_END, COL_W, 0.36,
        size=13, bold=True, color=CYAN)
    txt(slide, "For Learners", RX,  HEADER_END, COL_W, 0.36,
        size=13, bold=True, color=PINK)

    admin_steps = [
        ("Sign up & invite your team",
         "Create org account in minutes. Add users individually or bulk-upload a CSV."),
        ("Enable modules & set deadlines",
         "Choose which modules apply per department. Set completion due dates."),
        ("Track, export & certify",
         "Monitor compliance live. Export audit logs. Download certificates."),
    ]
    learner_steps = [
        ("Accept invite & set password",
         "Secure email link. Learner sets password. Logs in instantly."),
        ("Complete slides & pass quiz",
         "Self-paced slides then 10-question quiz — 70% pass mark, unlimited retakes."),
        ("Earn & share certificate",
         "Instant PDF certificate with QR code. Valid 12 months. Share on LinkedIn."),
    ]
    for i, (h_, b_) in enumerate(admin_steps):
        step(slide, i + 1, h_, b_, MX,  HEADER_END + 0.44 + i * 1.02, COL_W, accent=CYAN)
    for i, (h_, b_) in enumerate(learner_steps):
        step(slide, i + 1, h_, b_, RX,  HEADER_END + 0.44 + i * 1.02, COL_W, accent=PINK)

    # Centre arrow
    ax = MX + COL_W + 0.04
    txt(slide, "→", ax, HEADER_END + 1.20, 0.20, 0.40,
        size=18, bold=True, color=DGRAY, align=PP_ALIGN.CENTER)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 06  CERTIFICATES
# Left mock: x=MX  w=7.10  h=4.70
# Right features: x=8.00  w=4.88  6 rows at 0.78 pitch
# ══════════════════════════════════════════════════════════════════════════════
def s_certs(prs):
    CX, CW_, CH = MX, 7.10, 4.70    # cert mock
    FX = CX + CW_ + 0.25            # feature col x = 7.80
    FW = W - MX - FX                # 4.988"  x+w = 12.838 ✓

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "VERIFIABLE CERTIFICATES",
           "Every Completion is Proof-Ready",
           badge_color=PINK,
           subtitle="QR-verified PDF certificates — accepted by regulators, archived forever, shared on LinkedIn.")

    # ── Certificate mock ──────────────────────────────────────────────────────
    box(slide, CX, HEADER_END, CW_, CH, fill=SURF2)
    rule(slide, CX, HEADER_END, CW_, color=PINK, h=0.10)

    txt(slide, "CERTIFICATE OF COMPLETION",
        CX, HEADER_END + 0.18, CW_, 0.32,
        size=10, bold=True, color=PINK, align=PP_ALIGN.CENTER)

    txt(slide, "Rajesh Kumar",
        CX, HEADER_END + 0.60, CW_, 0.52,
        size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(slide, "has successfully completed",
        CX, HEADER_END + 1.18, CW_, 0.28,
        size=11, color=LGRAY, align=PP_ALIGN.CENTER)

    txt(slide, "POSH Act 2013\nPrevention of Sexual Harassment",
        CX, HEADER_END + 1.52, CW_, 0.72,
        size=15, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    txt(slide, "Score: 90%   |   Issued: 24 May 2026   |   Expires: 24 May 2027",
        CX, HEADER_END + 2.36, CW_, 0.30,
        size=9, color=GRAY, align=PP_ALIGN.CENTER)

    rule(slide, CX + 0.30, HEADER_END + 2.74, CW_ - 0.60, color=DGRAY, h=0.03)

    # QR placeholder box
    box(slide, CX + 5.35, HEADER_END + 2.95, 1.30, 1.35, fill=SURFACE)
    txt(slide, "[QR]", CX + 5.35, HEADER_END + 3.28, 1.30, 0.60,
        size=18, color=DGRAY, align=PP_ALIGN.CENTER)

    txt(slide, "Scan to verify:  training.aaplus.app/verify",
        CX + 0.30, HEADER_END + 4.30, CW_ - 1.80, 0.28,
        size=9, color=DGRAY, align=PP_ALIGN.LEFT)

    # ── Feature list ──────────────────────────────────────────────────────────
    features = [
        (PINK,   "Instant PDF",       "Issued automatically when learner passes"),
        (CYAN,   "QR Verification",   "Auditors scan to verify authenticity online"),
        (PURPLE, "Org Branding",      "Company logo on every certificate (Scale+)"),
        (GREEN,  "12-Month Validity", "Expiry shown; renewal reminders auto-sent"),
        (ORANGE, "Bulk Export",       "Admin bulk-downloads all certs as a ZIP file"),
        (VIOLET, "LinkedIn Share",    "One-click post to professional profile"),
    ]
    for i, (ac, title, body) in enumerate(features):
        fy = HEADER_END + i * 0.78   # last: 2.15 + 5*0.78 = 6.05  y+h=6.73 ✓
        box(slide, FX, fy, FW, 0.72, fill=SURFACE)
        rule(slide, FX, fy, FW, color=ac)
        txt(slide, title, FX + 0.13, fy + 0.09, FW - 0.26, 0.30,
            size=11, bold=True)
        txt(slide, body,  FX + 0.13, fy + 0.42, FW - 0.26, 0.24,
            size=9, color=LGRAY)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 07  ADMIN DASHBOARD
# Left col: stats row + 3 feature rows
# Right col: 5 feature rows
# col_w = (CW - 0.25) / 2 = 6.092"
# ══════════════════════════════════════════════════════════════════════════════
def s_admin(prs):
    COL_W = (CW - 0.25) / 2   # 6.092"
    RX    = MX + COL_W + 0.25  # 6.792"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "ADMIN DASHBOARD",
           "Complete Visibility. Zero Spreadsheets.",
           badge_color=CYAN,
           subtitle="Real-time compliance tracking, bulk management tools, and export-ready audit logs.")

    # ── Left: stat row ────────────────────────────────────────────────────────
    txt(slide, "Live Dashboard — What Admins See",
        MX, HEADER_END, COL_W, 0.35, size=12, bold=True, color=CYAN)

    mini = [
        ("73%", "Compliance",  CYAN),
        ("142", "Active",      PURPLE),
        ("18",  "Pending",     ORANGE),
        ("6",   "Overdue",     PINK),
    ]
    SW  = (COL_W - 3 * 0.08) / 4   # 1.473"
    SY  = HEADER_END + 0.44
    for i, (val, lbl, vc) in enumerate(mini):
        sx = MX + i * (SW + 0.08)
        box(slide, sx, SY, SW, 1.12, fill=SURFACE)
        rule(slide, sx, SY, SW, color=vc)
        txt(slide, val, sx, SY + 0.10, SW, 0.52,
            size=26, bold=True, color=vc, align=PP_ALIGN.CENTER)
        txt(slide, lbl, sx, SY + 0.68, SW, 0.30,
            size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # Feature rows left  (3 rows, 0.72 pitch, start at SY+1.12+0.12 = 3.83)
    left_feats = [
        (CYAN,   "Real-Time Dashboard",
         "Completion % across all departments and individuals"),
        (PURPLE, "Bulk User Management",
         "CSV invite, role changes, revoke access instantly"),
        (PINK,   "Module Control",
         "Enable per dept, set deadlines, lock modules (Corporate+)"),
    ]
    LFY = SY + 1.12 + 0.12   # 3.83
    for i, (ac, t, b) in enumerate(left_feats):
        fy = LFY + i * 0.72   # last: 3.83 + 1.44 = 5.27  y+h=5.93 ✓
        box(slide, MX, fy, COL_W, 0.66, fill=SURF2)
        rule(slide, MX, fy, COL_W, color=ac)
        txt(slide, t, MX + 0.13, fy + 0.08, COL_W - 0.26, 0.28,
            size=11, bold=True)
        txt(slide, b, MX + 0.13, fy + 0.38, COL_W - 0.26, 0.24,
            size=9, color=LGRAY)

    # ── Right: feature rows (5 rows, 0.82 pitch) ─────────────────────────────
    txt(slide, "Tools & Automations",
        RX, HEADER_END, COL_W, 0.35, size=12, bold=True, color=PINK)

    right_feats = [
        (ORANGE, "Excel Report Export",
         "Download compliance data + full audit trail as .xlsx"),
        (GREEN,  "Automated Email Nudges",
         "Reminders sent 7, 3, and 1 day before deadline"),
        (VIOLET, "Slack / Teams Webhooks",
         "Push notifications to team channels (Scale plan+)"),
        (CYAN,   "AI Policy Quiz",
         "Auto-assigns correct AI tier training module to all users"),
        (ORANGE, "POSH Wizard",
         "ICC details embedded live in POSH module slides"),
    ]
    RFY = HEADER_END + 0.44  # 2.59
    for i, (ac, t, b) in enumerate(right_feats):
        fy = RFY + i * 0.88   # last: 2.59 + 3.52 = 6.11  y+h=6.77 ✓
        box(slide, RX, fy, COL_W, 0.82, fill=SURFACE)
        rule(slide, RX, fy, COL_W, color=ac)
        txt(slide, t, RX + 0.13, fy + 0.09, COL_W - 0.26, 0.30,
            size=11, bold=True)
        txt(slide, b, RX + 0.13, fy + 0.43, COL_W - 0.26, 0.30,
            size=9, color=LGRAY)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 08  PRICING  (5 columns)
# CARD_W = 2.35  GAP = (CW - 5*2.35)/4 = 0.171
# Regular cards: cy=HEADER_END h=4.35
# Featured (Professional): cy=HEADER_END-0.15 h=4.50
# ══════════════════════════════════════════════════════════════════════════════
def s_pricing(prs):
    CARD_W = 2.35
    GAP    = (CW - 5 * CARD_W) / 4   # 0.171"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "PRICING",
           "Simple, Transparent Annual Pricing in INR",
           subtitle="All plans: 11+ modules + certificates + reports.  14-day free trial — no credit card.")

    tiers = [
        ("Essentials",   "Up to 25 users",  "Rs 18,000",  "Rs 720/user/yr",  PURPLE, False,
         ["All 11+ modules", "Digital certificates", "1 Admin account", "Basic compliance report"]),
        ("Business",     "Up to 100 users", "Rs 42,000",  "Rs 420/user/yr",  CYAN,   False,
         ["Everything in Essentials", "3 Admins + departments", "Bulk CSV invite",
          "Excel export & audit log", "Automated deadline nudges"]),
        ("Professional", "Up to 250 users", "Rs 80,000",  "Rs 320/user/yr",  PINK,   True,
         ["Everything in Business", "POSH Wizard + ICC details", "RPT Simulator",
          "AI Policy Tier module", "Org logo white-labeling", "Slack / Teams notifications"]),
        ("Corporate",    "Up to 500 users", "Rs 1,40,000","Rs 280/user/yr",  ORANGE, False,
         ["Everything in Professional", "Unlimited Admins", "Priority 24hr support",
          "Dedicated onboarding call", "2-year price lock"]),
        ("Enterprise",   "500+ employees",  "Custom",     "Rs 150+/user/yr", VIOLET, False,
         ["Custom modules + content", "Dedicated CSM + SLA", "Volume discount 15-25%",
          "Parent + child org mgmt"]),
    ]

    for i, (name, users, price, per_user, ac, featured, buls) in enumerate(tiers):
        cx  = MX + i * (CARD_W + GAP)
        cy  = HEADER_END - (0.15 if featured else 0)
        ch  = (FOOTER_Y - HEADER_END - 0.48) + (0.15 if featured else 0)   # ~4.32 / 4.47

        box(slide, cx, cy, CARD_W, ch,
            fill=SURF2 if featured else SURFACE,
            border=ac if featured else None)
        rule(slide, cx, cy, CARD_W, color=ac, h=0.08)

        if featured:
            chip(slide, "MOST POPULAR", cx + 0.20, cy - 0.30, CARD_W - 0.40, 0.24,
                 fill=ac, size=8)

        txt(slide, name,     cx + 0.12, cy + 0.12, CARD_W - 0.24, 0.32,
            size=13, bold=True)
        txt(slide, users,    cx + 0.12, cy + 0.44, CARD_W - 0.24, 0.26,
            size=9, color=GRAY)
        txt(slide, price,    cx + 0.12, cy + 0.75, CARD_W - 0.24, 0.52,
            size=19, bold=True, color=ac)
        txt(slide, per_user, cx + 0.12, cy + 1.26, CARD_W - 0.24, 0.26,
            size=9, color=LGRAY)
        rule(slide, cx + 0.12, cy + 1.56, CARD_W - 0.24, color=DGRAY, h=0.03)
        for j, b in enumerate(buls):
            txt(slide, f"+ {b}",
                cx + 0.12, cy + 1.66 + j * 0.44, CARD_W - 0.24, 0.40,
                size=9, color=LGRAY)

    # Insight bar
    IY = FOOTER_Y - 0.52
    box(slide, MX, IY, CW, 0.40, fill=SURF2)
    rule(slide, MX, IY, CW, color=PURPLE)
    txt(slide, "Rs 42,000 / 100 users / 365 days = Rs 1.15 per employee per day — less than a chai.",
        MX + 0.14, IY + 0.07, CW - 0.28, 0.26,
        size=11, italic=True, color=LGRAY)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 09  COMPETITIVE COMPARISON  (table)
# Cols: [Platform 2.45][India? 1.30][Regs 1.35][Certs 1.55][Cost 2.10][Note 3.28]
# Total used: 0.45 + 2.45+1.30+1.35+1.55+2.10+3.28 + 5*0.05 = 12.73 ✓
# ══════════════════════════════════════════════════════════════════════════════
def s_compare(prs):
    XS = [MX,   3.00, 4.40, 5.85, 7.50, 9.70]
    WS = [2.45, 1.30, 1.35, 1.55, 2.10, 3.28]
    # verify last col: 9.70 + 3.28 = 12.98  vs W-0.35=12.983 ✓
    HDR_H = 0.56
    ROW_H = 0.54
    GAP   = 0.04

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "COMPETITIVE COMPARISON",
           "AA Plus Wins on Value, Coverage, and Price",
           badge_color=GREEN,
           subtitle="Only platform combining India content + verifiable certificates + audit trail at this price.")

    headers = ["Platform", "India\nContent?", "Regs\nCovered",
               "Certs +\nAudit?", "100 Users\n(Annual)", "Note"]
    for i, (h_, x, w) in enumerate(zip(headers, XS, WS)):
        box(slide, x, HEADER_END, w - GAP, HDR_H, fill=SURF2)
        rule(slide, x, HEADER_END, w - GAP, color=GREEN)
        txt(slide, h_, x + 0.08, HEADER_END + 0.06, w - GAP - 0.16, HDR_H - 0.06,
            size=9, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)

    rows = [
        ("AA Plus (Business)", "Yes — 11+", "11+", "Yes + QR", "Rs 42,000",   "Best value for India compliance",    GREEN),
        ("POSH Check",         "POSH only",  "1",   "Basic",    "Rs 55,000",   "POSH only — no other regulations",   PINK),
        ("Ethena (US)",        "No",          "4-5", "Yes",      "~Rs 1,70,000","US-focused, no India law modules",   PINK),
        ("TalentLMS Grow",     "None",        "0",   "Basic",    "~Rs 2,30,000","LMS only — you build all content",   PINK),
        ("Disprz",             "Some",        "2-3", "Yes",      "~Rs 3,00,000","Full LMS but expensive for India",   PINK),
        ("KnowBe4",            "No",          "1",   "Yes",      "~Rs 3,40,000","Cybersecurity awareness only",       PINK),
        ("SAP Litmos",         "No",          "3-4", "Yes",      "Rs 10 L+",    "Enterprise MNC pricing",             PINK),
    ]
    DY = HEADER_END + HDR_H + GAP
    for ri, row_data in enumerate(rows):
        *vals, cost_color = row_data
        ry   = DY + ri * (ROW_H + GAP)
        aa   = ri == 0
        rbg  = SURF2 if aa else SURFACE
        for ci, (v, x, w) in enumerate(zip(vals, XS, WS)):
            box(slide, x, ry, w - GAP, ROW_H, fill=rbg)
            if aa:
                rule(slide, x, ry, w - GAP, color=GREEN)
            col = (GREEN if aa and ci == 4 else
                   PINK  if (not aa) and ci == 4 else
                   WHITE if aa else LGRAY)
            al = PP_ALIGN.CENTER if ci in (1, 2, 3) else PP_ALIGN.LEFT
            txt(slide, v, x + 0.08, ry + 0.10, w - GAP - 0.16, ROW_H - 0.14,
                size=9, bold=aa, color=col, align=al)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 10  RISK / ROI  (2 x 2 grid)
# CW2=6.08  CH2=2.26
# ══════════════════════════════════════════════════════════════════════════════
def s_roi(prs):
    CW2, CH2, GAP_Y = 6.08, 2.26, 0.14
    GAP_X = CW - 2 * CW2   # 0.273"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "THE RISK OF NON-COMPLIANCE",
           "The Cost of Getting It Wrong",
           badge_color=ORANGE,
           subtitle="Non-compliance with India regulations carries real financial, legal, and reputational consequences.")

    risks = [
        (PINK,   "POSH Act 2013",
         ["Penalty up to Rs 50,000 for first offence",
          "Rs 1 lakh+ for repeat or wilful default",
          "Criminal liability for the employer",
          "SEBI mandates documentary evidence"]),
        (CYAN,   "DPDP Act 2023",
         ["Penalties up to Rs 250 crore for failures",
          "Mandatory breach notification rules",
          "Board-level accountability under the Act",
          "Applies to all orgs handling personal data"]),
        (ORANGE, "CERT-In Cyber Directive",
         ["Incident reporting required within 6 hours",
          "Non-compliance invites MeitY audit",
          "Covers all companies with India operations",
          "Staff awareness training is documented req."]),
        (PURPLE, "SEBI / BSE Compliance",
         ["Annual board certification of POSH required",
          "ICC constitution must be documented",
          "Digital audit trail expected post-2023",
          "Non-compliance triggers exchange scrutiny"]),
    ]
    for i, (ac, title, buls) in enumerate(risks):
        col, row = i % 2, i // 2
        x = MX + col * (CW2 + GAP_X)
        y = HEADER_END + row * (CH2 + GAP_Y)
        info_card(slide, title, buls, x, y, CW2, CH2, accent=ac)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 11  WHO WE SERVE  (3 x 2 grid)
# COL_W=4.03  ROW_H=2.28
# ══════════════════════════════════════════════════════════════════════════════
def s_who(prs):
    GAP_X = 0.18
    COL_W = (CW - 2 * GAP_X) / 3   # 4.031"
    GAP_Y = 0.14
    ROW_H = (FOOTER_Y - HEADER_END - GAP_Y) / 2   # 2.33"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "WHO WE SERVE",
           "Every Indian Organisation That Has Employees",
           badge_color=VIOLET,
           subtitle="From 25-person startups to 5,000-employee conglomerates — all plans scale with you.")

    segments = [
        (PURPLE, "IT & Tech Startups",
         ["CERT-In cyber awareness (mandatory)", "DPDP data privacy for product teams",
          "Code of conduct for young teams"]),
        (CYAN,   "BFSI",
         ["Anti-corruption + FCPA compliance", "RPT Simulator for listed companies",
          "AI Policy Tier framework"]),
        (PINK,   "Manufacturing & FMCG",
         ["POSH + HSE modules bundled", "Whistleblower policy training",
          "Bulk CSV for large shift workforces"]),
        (ORANGE, "Professional Services",
         ["POSH mandatory for every firm", "AI Policy for consulting orgs",
          "White-label certs for client audits"]),
        (GREEN,  "Healthcare & Pharma",
         ["DPDP + Information Security", "Anti-corruption for pharma reps",
          "Audit trail for NABH / JCI"]),
        (VIOLET, "Conglomerates & Groups",
         ["Enterprise: parent + child org mgmt", "Volume pricing for 500+ users",
          "Custom SLA + dedicated CSM"]),
    ]
    for i, (ac, title, buls) in enumerate(segments):
        col, row = i % 3, i // 3
        x = MX + col * (COL_W + GAP_X)
        y = HEADER_END + row * (ROW_H + GAP_Y)
        info_card(slide, title, buls, x, y, COL_W, ROW_H, accent=ac)

    footer(slide)


# ══════════════════════════════════════════════════════════════════════════════
# 12  CTA
# ══════════════════════════════════════════════════════════════════════════════
def s_cta(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)

    # Top band  0..3.70
    box(slide, 0, 0, W, 3.70, fill=SURFACE)
    rule(slide, 0, 0, W, color=PURPLE, h=0.12)

    txt(slide, "Start Your 14-Day Free Trial",
        MX, 0.46, CW, 0.92, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(slide, "Full platform access  •  No credit card required  •  Cancel any time.",
        MX, 1.52, CW, 0.40, size=15, color=LGRAY, align=PP_ALIGN.CENTER)

    # URL chip  centred  w=5.20
    uw = 5.20
    chip(slide, "training.aaplus.app/signup",
         (W - uw) / 2, 2.16, uw, 0.48, fill=PURPLE, color=WHITE, size=14)

    # 3 action cards  CW3 = (CW - 2*0.20) / 3 = 4.011"
    CW3 = (CW - 2 * 0.20) / 3
    actions = [
        (PURPLE, "Free Trial",     "Sign up for 14 days full access.\nNo credit card. Cancel any time."),
        (CYAN,   "Book a Demo",    "30-minute live walkthrough for your\nHR and Legal team."),
        (PINK,   "POSH-Only Plan", "Start at Rs 25,000/year.\nAny org size. Any user count."),
    ]
    for i, (ac, title, body) in enumerate(actions):
        cx = MX + i * (CW3 + 0.20)
        box(slide, cx, 3.90, CW3, 2.60, fill=SURF2)
        rule(slide, cx, 3.90, CW3, color=ac, h=0.09)
        txt(slide, title, cx + 0.15, 4.06, CW3 - 0.30, 0.40,
            size=15, bold=True, color=ac)
        txt(slide, body, cx + 0.15, 4.56, CW3 - 0.30, 1.00,
            size=12, color=LGRAY)

    # Contact bar  6.65..7.5
    box(slide, 0, 6.65, W, H - 6.65, fill=SURFACE)
    txt(slide,
        "praveen@aaplusconsultants.com   |   training.aaplus.app   |   AA Plus Consultants Pvt. Ltd.",
        MX, 6.82, CW, 0.40, size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    s_cover(prs)
    s_problem(prs)
    s_solution(prs)
    s_modules(prs)
    s_how(prs)
    s_certs(prs)
    s_admin(prs)
    s_pricing(prs)
    s_compare(prs)
    s_roi(prs)
    s_who(prs)
    s_cta(prs)

    out = r"c:\Users\pc\Desktop\PolicyTraining\AA_Plus_Sales_Pitch_Deck.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
